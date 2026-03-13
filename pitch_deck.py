"""Generate OpenClaw AI Lead Intelligence Pitch Deck"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

GREEN = RGBColor(0x2D, 0x8B, 0x2D)
DARK_GREEN = RGBColor(0x1A, 0x5C, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
RED = RGBColor(0xCC, 0x33, 0x33)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return txBox


def add_bullets(slide, items, left, top, width, size=16, color=BLACK, spacing=10):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(spacing)


# ========== SLIDE 1: TITLE ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_GREEN)

add_text(slide, Inches(1), Inches(1.2), Inches(11), Inches(1.5),
         "OpenClaw", 60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(2.8), Inches(11), Inches(1),
         "AI-Powered Lead Intelligence & Personalized Outreach Agent", 28, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.6),
         "Prospect Discovery  •  Digital Footprint Analysis  •  Lead Scoring  •  Personalized Messaging",
         18, color=RGBColor(0xA0, 0xD0, 0xA0), align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
         "SUSS OpenClaw Project", 16, color=RGBColor(0x80, 0xB0, 0x80), align=PP_ALIGN.CENTER)


# ========== SLIDE 2: THE CASE ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.2), DARK_GREEN)
add_text(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.7),
         "THE CASE", 36, bold=True, color=WHITE)

add_text(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(1),
         "An intelligent agent that analyzes a prospect's digital footprint and automatically\ngenerates high-quality personalized outreach messages.", 22, bold=True, color=BLACK)

add_text(slide, Inches(0.8), Inches(3), Inches(11), Inches(0.6),
         "Instead of sending generic cold emails, businesses can engage leads with\ncontext-aware messaging at scale.", 20, color=GRAY)

# Problem stats
problems = [
    "❌  97% of cold outreach emails are ignored (generic, irrelevant, impersonal)",
    "❌  Sales teams spend 40% of time researching prospects manually",
    "❌  Traditional CRMs store data but don't generate intelligence",
    "❌  No tool connects prospect research → tone analysis → message generation",
]
add_bullets(slide, problems, Inches(0.8), Inches(4.2), Inches(11), size=17, color=BLACK)


# ========== SLIDE 3: SYSTEM FLOW ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.2), DARK_GREEN)
add_text(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.7),
         "SYSTEM FLOW", 36, bold=True, color=WHITE)

steps = [
    ("1", "Prospect\nDiscovery",
     "• LinkedIn profile search\n• Company websites\n• Hiring announcements\n• Industry news"),
    ("2", "Digital Footprint\nAnalysis",
     "• Text signals (social posts, articles)\n• Structured data (company size, role)\n• Visual signals (marketing content)\n• Online activity patterns"),
    ("3", "Lead\nScoring",
     "• Priority P1-P5 ranking\n• Role-product fit analysis\n• Company relevance scoring\n• Decision-maker identification"),
    ("4", "Personalized\nMessaging",
     "• Tone-matched connection msgs\n• Customized follow-ups\n• Meeting talking points\n• Multi-channel outreach"),
]

for i, (num, title, desc) in enumerate(steps):
    left = Inches(0.4 + i * 3.2)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(1.1), Inches(1.5), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = GREEN
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Arrow (except last)
    if i < 3:
        add_text(slide, left + Inches(2.8), Inches(1.55), Inches(0.5), Inches(0.6),
                 "→", 28, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    add_text(slide, left, Inches(2.4), Inches(3), Inches(0.8),
             title, 18, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.2), Inches(3.4), Inches(2.6), Inches(3),
             desc, 14, color=GRAY, align=PP_ALIGN.LEFT)


# ========== SLIDE 4: DIGITAL FOOTPRINT DEEP DIVE ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.2), DARK_GREEN)
add_text(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.7),
         "DIGITAL FOOTPRINT ANALYSIS", 36, bold=True, color=WHITE)

add_text(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
         "Signals Into Meaning — We analyze 4 types of digital signals per prospect", 22, bold=True, color=BLACK)

signals = [
    ("Text Signals", "Social media posts, LinkedIn articles,\nblog content, website bios,\nconference talks, quotes in press",
     "→ Communication style, vocabulary,\n   topics they care about, tone"),
    ("Structured Data", "Company size, industry, revenue,\nhiring trends, funding rounds,\njob title & seniority level",
     "→ Company fit score, budget likelihood,\n   decision-making authority"),
    ("Visual Signals", "Product screenshots, marketing\ncreatives, brand positioning,\ncompany culture imagery",
     "→ Brand personality, company maturity,\n   visual communication preferences"),
    ("Activity Signals", "Website visits, engagement patterns,\nevent attendance, group memberships,\ncontent sharing behavior",
     "→ Interest level, buying intent,\n   preferred communication channels"),
]

for i, (title, desc, output) in enumerate(signals):
    top = Inches(2.3 + i * 1.2)
    add_shape(slide, Inches(0.8), top, Inches(2.5), Inches(1), LIGHT_GRAY)
    add_text(slide, Inches(1), top + Inches(0.05), Inches(2.1), Inches(0.4),
             title, 16, bold=True, color=GREEN)
    add_text(slide, Inches(3.6), top + Inches(0.05), Inches(4.5), Inches(0.9),
             desc, 13, color=BLACK)
    add_text(slide, Inches(8.5), top + Inches(0.05), Inches(4.5), Inches(0.9),
             output, 13, bold=True, color=DARK_GREEN)


# ========== SLIDE 5: LEAD SCORING ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.2), DARK_GREEN)
add_text(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.7),
         "AI LEAD SCORING", 36, bold=True, color=WHITE)

add_text(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
         "Likelihood of lead to convert — scored P1 (low) to P5 (high)", 22, bold=True, color=BLACK)

scores = [
    ("P5", GREEN, "Direct decision maker for your product category\nStrong role-product fit, active buyer signals"),
    ("P4", RGBColor(0x4A, 0xA8, 0x4A), "Senior role at target company\nRelevant department, likely budget authority"),
    ("P3", RGBColor(0xCC, 0xA0, 0x00), "Relevant role but mid-level or smaller company\nMay need internal champion to proceed"),
    ("P2", RGBColor(0xCC, 0x80, 0x33), "Tangentially related role\nLong-shot but possible with right angle"),
    ("P1", RGBColor(0xCC, 0x55, 0x55), "Wrong function, too junior, or poor fit\nDeprioritize or skip"),
]

for i, (label, color, desc) in enumerate(scores):
    top = Inches(2.3 + i * 0.95)
    add_shape(slide, Inches(0.8), top, Inches(1), Inches(0.75), color)
    add_text(slide, Inches(0.8), top + Inches(0.1), Inches(1), Inches(0.5),
             label, 24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(2.2), top + Inches(0.1), Inches(10), Inches(0.6),
             desc, 15, color=BLACK)


# ========== SLIDE 6: PERSONALIZED MESSAGING ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.2), DARK_GREEN)
add_text(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.7),
         "PERSONALIZED MESSAGING", 36, bold=True, color=WHITE)

add_text(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
         "AI generates outreach that mirrors each prospect's communication style", 22, bold=True, color=BLACK)

# Generic vs Personalized comparison
add_shape(slide, Inches(0.5), Inches(2.3), Inches(5.8), Inches(4.5), LIGHT_GRAY)
add_text(slide, Inches(0.8), Inches(2.5), Inches(5.2), Inches(0.5),
         "❌  GENERIC OUTREACH", 20, bold=True, color=RED)
add_text(slide, Inches(0.8), Inches(3.2), Inches(5.2), Inches(3.5),
         '"Hi [Name],\n\nI wanted to reach out because our company\noffers solutions that might be relevant to\nyour business.\n\nWould you be open to a quick call?"\n\n→ 2-3% response rate\n→ Feels like spam\n→ No personalization', 15, color=GRAY)

add_shape(slide, Inches(7), Inches(2.3), Inches(5.8), Inches(4.5), LIGHT_GRAY)
add_text(slide, Inches(7.3), Inches(2.5), Inches(5.2), Inches(0.5),
         "✓  OPENCLAW OUTREACH", 20, bold=True, color=GREEN)
add_text(slide, Inches(7.3), Inches(3.2), Inches(5.2), Inches(3.5),
         '"Hi Badai, your work driving employee\nbenefits at Reckitt Indonesia caught my\nattention — especially the wellness\nprograms across 2,000+ staff.\n\nWe help HR leaders like you quantify\nthe ROI of wellness investments..."\n\n→ 15-25% response rate\n→ Feels like a warm intro\n→ Tone-matched to their style', 15, color=BLACK)


# ========== SLIDE 7: LIVE DASHBOARD ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.2), DARK_GREEN)
add_text(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.7),
         "TRACKING DASHBOARD", 36, bold=True, color=WHITE)

add_text(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
         "Real-time pipeline management with built-in contact tracking", 22, bold=True, color=BLACK)

features = [
    "📊  Priority-sorted lead cards (P1-P5) with one-click LinkedIn profile access",
    "📋  Contact status tracking: Not Contacted → Contacted → Replied → Meeting Scheduled",
    "🔍  Filter by priority score, contact status, or role type (HR / CFO)",
    "📝  Per-lead notes, contact dates, and follow-up reminders",
    "👤  Google background intelligence on each prospect's character and interests",
    "💬  Pre-written outreach messages ready to copy-paste into LinkedIn",
    "📈  Pipeline metrics: total leads, contacted %, reply rate, meetings booked",
    "📥  One-click export to Excel with all tracking data",
    "🔄  Hourly auto-refresh to keep data current",
]
add_bullets(slide, features, Inches(0.8), Inches(2.4), Inches(11), size=17, color=BLACK, spacing=8)


# ========== SLIDE 8: BUSINESS VALUE ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.2), DARK_GREEN)
add_text(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.7),
         "BUSINESS VALUE", 36, bold=True, color=WHITE)

values = [
    ("Higher Response Rate", "15-25%",
     "Tone-matched, context-aware messages\nget 5-10x more replies than generic outreach"),
    ("Better Targeting", "P1-P5",
     "AI scoring prioritizes leads most likely\nto convert — stop wasting time on bad fits"),
    ("Scalable Personalization", "1000s",
     "Reach thousands of prospects simultaneously\nwith individually crafted messages"),
    ("Time Savings", "80%",
     "Automate prospect research, background\nchecks, and message drafting — hours → minutes"),
]

for i, (title, metric, desc) in enumerate(values):
    left = Inches(0.5 + i * 3.2)
    add_shape(slide, left, Inches(1.5), Inches(2.8), Inches(5), LIGHT_GRAY)
    add_text(slide, left, Inches(1.7), Inches(2.8), Inches(0.5),
             title, 18, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(slide, left, Inches(2.4), Inches(2.8), Inches(1),
             metric, 48, bold=True, color=DARK_GREEN, align=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.2), Inches(3.8), Inches(2.4), Inches(2),
             desc, 14, color=GRAY, align=PP_ALIGN.CENTER)


# ========== SLIDE 9: TECH STACK & ARCHITECTURE ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.2), DARK_GREEN)
add_text(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.7),
         "TECHNOLOGY", 36, bold=True, color=WHITE)

add_text(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
         "Built with Python • Powered by AI • Deployed on OpenClaw", 22, bold=True, color=BLACK)

stack = [
    ("Data Collection", "Python + Firecrawl\nWeb search, LinkedIn discovery\nGoogle background checks\nSocial media scraping"),
    ("AI Analysis", "Claude AI (via Claude Code)\nTone & style profiling\nPersonality analysis\nMessage generation"),
    ("Dashboard", "Streamlit\nReal-time tracking\nPipeline management\nExcel export"),
    ("Deployment", "OpenClaw Platform\nScheduled runs\nAuto-refresh\nMulti-user access"),
]

for i, (title, desc) in enumerate(stack):
    left = Inches(0.5 + i * 3.2)
    add_shape(slide, left, Inches(2.3), Inches(2.8), Inches(3.5), LIGHT_GRAY)
    add_text(slide, left, Inches(2.5), Inches(2.8), Inches(0.5),
             title, 18, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.2), Inches(3.2), Inches(2.4), Inches(2.5),
             desc, 15, color=GRAY, align=PP_ALIGN.CENTER)


# ========== SLIDE 10: DEMO RESULTS ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.2), DARK_GREEN)
add_text(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.7),
         "LIVE DEMO RESULTS", 36, bold=True, color=WHITE)

add_text(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
         "Case Study: Heyva Health — Selling to HR Directors & CFOs in Indonesia", 22, bold=True, color=BLACK)

metrics = [
    ("38", "Leads\nDiscovered"),
    ("12", "P4-P5\nHigh Priority"),
    ("38", "Personalized\nMessages"),
    ("12", "Google Background\nChecks"),
]
for i, (num, label) in enumerate(metrics):
    left = Inches(0.5 + i * 3.2)
    add_shape(slide, left, Inches(2.3), Inches(2.8), Inches(2), LIGHT_GRAY)
    add_text(slide, left, Inches(2.5), Inches(2.8), Inches(0.8),
             num, 48, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(slide, left, Inches(3.5), Inches(2.8), Inches(0.6),
             label, 15, color=GRAY, align=PP_ALIGN.CENTER)

highlights = [
    "✓  Discovered HR Directors at Reckitt Indonesia, Prudential, AXA Mandiri automatically",
    "✓  Each lead profiled with Google background check revealing character and interests",
    "✓  AI-generated outreach messages matched to each person's communication tone",
    "✓  Full pipeline from discovery → analysis → scoring → messaging in under 10 minutes",
]
add_bullets(slide, highlights, Inches(0.8), Inches(4.8), Inches(11), size=17, color=BLACK, spacing=6)


# ========== SLIDE 11: CTA ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_GREEN)

add_text(slide, Inches(1), Inches(1), Inches(11), Inches(1.2),
         "OpenClaw", 60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.2),
         "Stop sending cold emails.\nStart having warm conversations.", 28, color=WHITE, align=PP_ALIGN.CENTER)

add_shape(slide, Inches(3.5), Inches(4.5), Inches(6.333), Inches(1), GREEN)
add_text(slide, Inches(3.5), Inches(4.55), Inches(6.333), Inches(0.9),
         "See the Live Dashboard Demo →", 26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(6), Inches(11), Inches(0.8),
         "AI-Powered Prospect Discovery  •  Digital Footprint Analysis  •  Lead Scoring  •  Personalized Outreach",
         16, color=RGBColor(0xA0, 0xD0, 0xA0), align=PP_ALIGN.CENTER)


# Save
output_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "output")
os.makedirs(output_dir, exist_ok=True)
output_path = os.path.join(output_dir, "OpenClaw_Pitch_Deck.pptx")
prs.save(output_path)
print(f"Pitch deck saved to: {output_path}")
